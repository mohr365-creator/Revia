#!/usr/bin/env python3
"""Create a Revia-branded .pptx with the background pre-applied on every slide."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BACKGROUND_IMG = "/home/user/Revia/revia-background-vivid.png"
OUTPUT = "/home/user/Revia/revia-template.pptx"

CREAM  = RGBColor(0xFF, 0xF4, 0xE1)  # #fff4e1
AMBER  = RGBColor(0xE8, 0x95, 0x56)  # #e89556

W = Inches(13.333)   # 16:9 widescreen width
H = Inches(7.5)      # 16:9 widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def add_background(slide):
    """Stretch the background image to fill the slide."""
    pic = slide.shapes.add_picture(BACKGROUND_IMG, 0, 0, W, H)
    # Move picture to back (z-order index 0)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def add_text(slide, text, left, top, width, height,
             font_size=28, bold=False, color=CREAM, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Inter"


# ── Slide 1: Title slide ──────────────────────────────────────────────────────
layout = prs.slide_layouts[6]   # blank
slide = prs.slides.add_slide(layout)
add_background(slide)

add_text(slide, "REVIA",
         left=Inches(1), top=Inches(2.8), width=Inches(11.3), height=Inches(1.2),
         font_size=72, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

add_text(slide, "Restoring the Routes That Connect Us",
         left=Inches(1), top=Inches(4.2), width=Inches(11.3), height=Inches(0.7),
         font_size=24, color=CREAM, align=PP_ALIGN.CENTER)


# ── Slide 2: Section / content starter ───────────────────────────────────────
slide2 = prs.slides.add_slide(layout)
add_background(slide2)

add_text(slide2, "Slide Title",
         left=Inches(0.8), top=Inches(0.5), width=Inches(10), height=Inches(0.9),
         font_size=40, bold=True, color=AMBER)

add_text(slide2, "Add your content here.",
         left=Inches(0.8), top=Inches(1.6), width=Inches(11), height=Inches(5),
         font_size=22, color=CREAM)


prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
